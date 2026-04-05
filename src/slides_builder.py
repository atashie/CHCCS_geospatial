#!/usr/bin/env python3
"""
CHCCS Methods Overview -- Slide Deck Builder

Captures map screenshots via Playwright and assembles a .pptx deck
illustrating the two core methodological pipelines:
  1. Dasymetric downscaling  (Census -> parcels)
  2. Network travel-time proximity  (roads -> school communities)

Usage:
    python src/slides_builder.py                    # full run
    python src/slides_builder.py --skip-screenshots # reuse cached PNGs
    python src/slides_builder.py --output FILE      # custom output path
"""

import argparse
from pathlib import Path

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paths ────────────────────────────────────────────────────────────────────

PROJECT_ROOT = Path(__file__).resolve().parent.parent
MAPS_DIR = PROJECT_ROOT / "assets" / "maps"
SLIDES_DIR = PROJECT_ROOT / "assets" / "slides"
SCREENSHOTS_DIR = SLIDES_DIR / "screenshots"
DEFAULT_OUTPUT = SLIDES_DIR / "methods_overview.pptx"

# ── Screenshot settings ──────────────────────────────────────────────────────

VIEWPORT_W, VIEWPORT_H = 1280, 960  # 4:3 for slide embedding
TILE_WAIT_MS = 3000  # ms to wait after layer change for tile/render settling
MAP_CENTER = [35.9132, -79.0558]
MAP_ZOOM = 12

# ── JS helpers ───────────────────────────────────────────────────────────────

JS_FIND_MAP = """
() => {
    for (const key of Object.keys(window)) {
        try {
            const v = window[key];
            if (v && typeof v === 'object'
                && typeof v.getZoom === 'function'
                && typeof v.invalidateSize === 'function') {
                return key;
            }
        } catch(e) {}
    }
    return null;
}
"""

JS_HIDE_ELEMENTS = """
(selectors) => {
    for (const sel of selectors) {
        const el = document.querySelector(sel);
        if (el) el.style.display = 'none';
    }
}
"""

JS_RESIZE_MAP = """
(args) => {
    const [mapKey, lat, lng, zoom] = args;
    const m = window[mapKey];
    if (!m) return;
    m.invalidateSize();
    if (lat != null) m.setView([lat, lng], zoom);
}
"""


# ── Screenshot capture ───────────────────────────────────────────────────────

def _file_url(path):
    """Local path -> file:// URL."""
    return "file:///" + path.resolve().as_posix()


def _wait(page, ms=TILE_WAIT_MS):
    """Wait for network + extra settle time for canvas rendering."""
    try:
        page.wait_for_load_state("networkidle", timeout=15000)
    except Exception:
        # Tile servers may keep connections open; fall back to domcontentloaded
        try:
            page.wait_for_load_state("domcontentloaded", timeout=10000)
        except Exception:
            pass
    page.wait_for_timeout(ms)


def _hide_leaflet_chrome(page):
    """Hide standard Leaflet zoom/scale controls."""
    page.evaluate(JS_HIDE_ELEMENTS, [
        ".leaflet-control-zoom",
        ".leaflet-control-scale",
        ".leaflet-control-attribution",
    ])


def screenshot_socioeconomic(page, out_dir):
    """Capture three dasymetric-pipeline screenshots from the socioeconomic map.

    Pipeline:
      Block Groups choropleth (bg) -> Blocks choropleth (blk) -> Population Dots (dots)

    Steps 1-2 use Median Household Income (metric 1) so the choropleth polygons
    are visible at both resolutions.  Step 3 switches to Race/Ethnicity (metric 0)
    to show the parcel-constrained dot-density layer.
    """
    print("Capturing socioeconomic map screenshots...")
    page.goto(_file_url(MAPS_DIR / "school_socioeconomic_map.html"))
    _wait(page, ms=5000)  # initial load

    # Elements to hide for clean screenshots
    hide = [
        "#socio-banner", "#ctrl-panel", "#zone-strip",
        ".faq-btn", "#faq-panel",
    ]

    # All custom JS functions (updateMetric, updateLayerToggle) live inside a
    # DOMContentLoaded closure, so they are NOT globally accessible.
    # Instead we set radio-button state and dispatch 'change' events, which
    # trigger the addEventListener handlers that call the internal functions.

    # Step 1: Block Groups choropleth  (metric 1 = Median HH Income, layer "bg")
    page.evaluate("""
        () => {
            var mr = document.querySelector('input[name="metric"][value="1"]');
            var lr = document.querySelector('input[name="layer"][value="bg"]');
            if (lr) lr.checked = true;       // set layer BEFORE metric dispatch
            if (mr) { mr.checked = true; mr.dispatchEvent(new Event('change')); }
        }
    """)
    _wait(page, ms=4000)
    page.evaluate(JS_HIDE_ELEMENTS, hide)
    _hide_leaflet_chrome(page)
    page.screenshot(path=str(out_dir / "step1_block_groups.png"))
    print("  > step1_block_groups.png")

    # Step 2: Blocks (est.) choropleth  (same metric, layer "blk")
    page.evaluate("""
        () => {
            var lr = document.querySelector('input[name="layer"][value="blk"]');
            if (lr) { lr.checked = true; lr.dispatchEvent(new Event('change')); }
        }
    """)
    _wait(page, ms=4000)
    page.evaluate(JS_HIDE_ELEMENTS, hide)
    _hide_leaflet_chrome(page)
    page.screenshot(path=str(out_dir / "step2_blocks_est.png"))
    print("  > step2_blocks_est.png")

    # Step 3: Population Dots  (metric 0 = Race/Ethnicity, layer "dots")
    page.evaluate("""
        () => {
            var mr = document.querySelector('input[name="metric"][value="0"]');
            var lr = document.querySelector('input[name="layer"][value="dots"]');
            if (lr) lr.checked = true;       // set layer BEFORE metric dispatch
            if (mr) { mr.checked = true; mr.dispatchEvent(new Event('change')); }
        }
    """)
    _wait(page, ms=4000)
    page.evaluate(JS_HIDE_ELEMENTS, hide)
    _hide_leaflet_chrome(page)
    page.screenshot(path=str(out_dir / "step3_pop_dots.png"))
    print("  > step3_pop_dots.png")


def screenshot_community(page, out_dir):
    """Capture three travel-time-pipeline screenshots from the community map.

    Step progression:
      Road network -> Travel-time heatmap -> Heatmap + residential parcels
    """
    print("Capturing community map screenshots...")
    page.goto(_file_url(MAPS_DIR / "school_community_map.html"))
    _wait(page, ms=6000)  # longer: dynamic content

    # Wait for scenario radio buttons to be populated
    try:
        page.wait_for_selector("input[name=\"scenario\"]", timeout=15000)
    except Exception:
        print("  ! Warning: scenario controls not found; proceeding anyway")
    page.wait_for_timeout(1000)

    # Find map, hide sidebar, resize to fill viewport
    map_key = page.evaluate(JS_FIND_MAP)

    page.evaluate(JS_HIDE_ELEMENTS, [
        "#community-controls", "#affected-count-label", "#walkzone-count-label",
    ])
    _hide_leaflet_chrome(page)

    if map_key:
        page.evaluate(JS_RESIZE_MAP,
                      [map_key, MAP_CENTER[0], MAP_CENTER[1], MAP_ZOOM])
        page.wait_for_timeout(2000)

    # Step 1: Road network visible, other overlays off
    page.evaluate("""
        () => {
            var net = document.getElementById('show-network');
            var wz  = document.getElementById('show-walk-zones');
            var pp  = document.getElementById('show-properties');
            if (net) net.checked = true;
            if (wz)  wz.checked = false;
            if (pp)  pp.checked = false;
            if (window.updateCommunityMap) window.updateCommunityMap();
        }
    """)
    _wait(page, ms=4000)
    page.screenshot(path=str(out_dir / "step1_road_network.png"))
    print("  > step1_road_network.png")

    # Step 2: Travel-time heatmap only (hide road network for clarity)
    page.evaluate("""
        () => {
            var net = document.getElementById('show-network');
            if (net) net.checked = false;
            var abs = document.querySelector('input[name="layer_type"][value="abs"]');
            if (abs) { abs.checked = true; }
            if (window.updateCommunityMap) window.updateCommunityMap();
        }
    """)
    _wait(page, ms=4000)
    page.screenshot(path=str(out_dir / "step2_travel_time.png"))
    print("  > step2_travel_time.png")

    # Step 3: School zones + travel time from the closure analysis map.
    # This shows school attendance zone boundaries overlaid on the travel-time
    # gradient, illustrating how communities are defined by actual proximity.
    print("  Switching to closure analysis map for step 3...")
    page.goto(_file_url(MAPS_DIR / "school_closure_analysis.html"))
    _wait(page, ms=6000)

    # Switch to Part 1 (Travel Time) tab
    page.evaluate("if (window.switchTab) window.switchTab('part1');")
    page.wait_for_timeout(2000)

    # Configure: absolute travel time, zone boundaries ON, road network OFF
    page.evaluate("""
        () => {
            var absR = document.querySelector('input[name="p1-view"][value="abs"]');
            if (absR) { absR.checked = true; }
            var zones = document.getElementById('p1-show-zones');
            if (zones) zones.checked = true;
            var net = document.getElementById('p1-show-network');
            if (net) net.checked = false;
            if (window.updatePart1) window.updatePart1();
        }
    """)
    _wait(page, ms=4000)

    # Find map, hide sidebar, resize
    map_key2 = page.evaluate(JS_FIND_MAP)
    page.evaluate(JS_HIDE_ELEMENTS, ["#closure-controls", "#closure-banner",
                                      "#faq-panel-closure", ".faq-btn-closure"])
    _hide_leaflet_chrome(page)
    if map_key2:
        page.evaluate(JS_RESIZE_MAP,
                      [map_key2, MAP_CENTER[0], MAP_CENTER[1], MAP_ZOOM])
        page.wait_for_timeout(2000)

    page.screenshot(path=str(out_dir / "step3_communities.png"))
    print("  > step3_communities.png")


def capture_all_screenshots(out_dir):
    """Run Playwright to capture all map screenshots."""
    out_dir.mkdir(parents=True, exist_ok=True)

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        ctx = browser.new_context(viewport={"width": VIEWPORT_W, "height": VIEWPORT_H})
        page = ctx.new_page()

        screenshot_socioeconomic(page, out_dir)
        screenshot_community(page, out_dir)

        browser.close()

    print(f"All screenshots saved to {out_dir}")


# ── Slide assembly ───────────────────────────────────────────────────────────

# Design tokens
CLR_DARK  = RGBColor(0x2C, 0x3E, 0x50)
CLR_LIGHT = RGBColor(0xEC, 0xF0, 0xF1)
CLR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CLR_MID   = RGBColor(0x7F, 0x8C, 0x8D)
CLR_ACCENT = RGBColor(0x29, 0x80, 0xB9)


def _set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, left, top, width, height, text, *,
              size=14, bold=False, color=CLR_DARK, align=PP_ALIGN.LEFT,
              word_wrap=True):
    """Add a simple text box and return its paragraph for further tweaks."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = word_wrap
    par = tf.paragraphs[0]
    par.text = text
    par.font.size = Pt(size)
    par.font.bold = bold
    par.font.color.rgb = color
    par.alignment = align
    return tf


def slide_title(prs):
    """Slide 1 -- Title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _set_bg(slide, CLR_DARK)

    _add_text(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.2),
              "CHCCS Geospatial Analysis",
              size=44, bold=True, color=CLR_WHITE, align=PP_ALIGN.CENTER)

    _add_text(slide, Inches(1), Inches(3.3), Inches(11.3), Inches(0.8),
              "Methods Overview",
              size=28, color=RGBColor(0xBD, 0xC3, 0xC7), align=PP_ALIGN.CENTER)

    _add_text(slide, Inches(1), Inches(5.5), Inches(11.3), Inches(0.5),
              "Chapel Hill-Carrboro City Schools",
              size=18, color=RGBColor(0x95, 0xA5, 0xA6), align=PP_ALIGN.CENTER)


def slide_pipeline(prs, *, title, images, labels, caption, source):
    """Generic pipeline slide: three images with -> arrows between them."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    _add_text(slide, Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7),
              title, size=28, bold=True, color=CLR_DARK)

    # Layout arithmetic
    img_w = Inches(3.5)
    img_h = Inches(2.625)    # 4:3
    arrow_gap = Inches(0.6)
    total_w_in = 3.5 * 3 + 0.6 * 2  # 11.7"
    margin_left = (13.333 - total_w_in) / 2
    img_top = Inches(1.15)

    for i, (img_path, label_text) in enumerate(zip(images, labels)):
        x = Inches(margin_left + i * (3.5 + 0.6))

        # Image (with border placeholder if missing)
        if img_path.exists():
            slide.shapes.add_picture(str(img_path), x, img_top, img_w, img_h)
        else:
            # Gray placeholder box
            from pptx.enum.shapes import MSO_SHAPE
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, img_top, img_w, img_h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = CLR_LIGHT
            shape.line.color.rgb = CLR_MID

        # Label below image
        _add_text(slide, x, img_top + img_h + Inches(0.08), img_w, Inches(0.55),
                  label_text, size=13, bold=True, color=CLR_DARK,
                  align=PP_ALIGN.CENTER)

        # Arrow between images 1->2 and 2->3
        if i < 2:
            arr_x = x + img_w + Inches(0.02)
            arr_y = img_top + img_h / 2 - Inches(0.25)
            _add_text(slide, arr_x, arr_y, arrow_gap - Inches(0.04), Inches(0.5),
                      "\u2192", size=36, color=CLR_ACCENT, align=PP_ALIGN.CENTER)

    # Caption
    _add_text(slide, Inches(0.8), Inches(4.55), Inches(11.7), Inches(1.2),
              caption, size=13, color=CLR_MID)

    # Source line
    _add_text(slide, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
              source, size=10, color=CLR_MID, align=PP_ALIGN.RIGHT)


def slide_synthesis(prs, screenshots_dir):
    """Slide 4 -- Bringing It Together: demographics + proximity -> communities."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ss = screenshots_dir

    _add_text(slide, Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7),
              "Bringing It Together: People + Proximity",
              size=28, bold=True, color=CLR_DARK)

    # Two small input images stacked on left
    sm_w, sm_h = Inches(3.2), Inches(2.4)
    left_x = Inches(0.6)

    demo_img = ss / "step3_pop_dots.png"
    travel_img = ss / "step2_travel_time.png"
    combined_img = ss / "step3_communities.png"

    # Demographics thumbnail
    if demo_img.exists():
        slide.shapes.add_picture(str(demo_img), left_x, Inches(1.1), sm_w, sm_h)
    _add_text(slide, left_x, Inches(3.55), sm_w, Inches(0.3),
              "Who lives where (dasymetric)", size=11, bold=True,
              color=CLR_DARK, align=PP_ALIGN.CENTER)

    # Travel-time thumbnail
    if travel_img.exists():
        slide.shapes.add_picture(str(travel_img), left_x, Inches(4.0), sm_w, sm_h)
    _add_text(slide, left_x, Inches(6.45), sm_w, Inches(0.3),
              "How far is each school (network)", size=11, bold=True,
              color=CLR_DARK, align=PP_ALIGN.CENTER)

    # Merging arrow
    _add_text(slide, Inches(4.1), Inches(3.4), Inches(0.8), Inches(0.6),
              "\u2192", size=48, color=CLR_ACCENT, align=PP_ALIGN.CENTER)

    # Combined result (larger)
    big_w, big_h = Inches(7.8), Inches(5.85)
    if combined_img.exists():
        slide.shapes.add_picture(str(combined_img), Inches(5.1), Inches(0.95),
                                 big_w, big_h)

    _add_text(slide, Inches(5.1), Inches(6.9), big_w, Inches(0.4),
              "School community demographics: high-resolution population + actual travel proximity",
              size=11, color=CLR_MID, align=PP_ALIGN.CENTER)


def slide_summary(prs):
    """Slide 5 -- Methods Summary."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_text(slide, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.8),
              "Standard Methods, Open Data",
              size=32, bold=True, color=CLR_DARK)

    bullets = [
        ("Dasymetric mapping",
         "Census demographics refined to parcel level using residential zoning data"),
        ("Network travel-time analysis",
         "Shortest-path routing on real road networks with speed limits, "
         "traffic signals, and bike infrastructure"),
        ("Open, public data sources",
         "U.S. Census ACS & Decennial, OpenStreetMap, NCDOT traffic counts, "
         "FEMA flood zones"),
        ("All 11 schools treated equally",
         "Same methodology applied uniformly -- no assumptions about any school"),
        ("Transparent and replicable",
         "All code and data sources documented; results independently verifiable"),
    ]

    y_in = 1.5
    for heading, detail in bullets:
        _add_text(slide, Inches(1.5), Inches(y_in), Inches(10.3), Inches(0.4),
                  heading, size=20, bold=True, color=CLR_ACCENT)
        _add_text(slide, Inches(1.5), Inches(y_in + 0.38), Inches(10.3), Inches(0.5),
                  detail, size=14, color=CLR_MID)
        y_in += 1.0


def build_presentation(screenshots_dir, output_path):
    """Assemble the 5-slide .pptx deck."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    ss = screenshots_dir

    # Slide 1: Title
    slide_title(prs)

    # Slide 2: Dasymetric pipeline
    slide_pipeline(
        prs,
        title="Who Lives Where: From Census to Communities",
        images=[
            ss / "step1_block_groups.png",
            ss / "step2_blocks_est.png",
            ss / "step3_pop_dots.png",
        ],
        labels=[
            "Census Block Groups\n(~1,500 people each)",
            "Block-Level Downscaling\n(Census blocks + parcels)",
            "Parcel-Constrained Population\n(dasymetric method)",
        ],
        caption=(
            "Census demographics (ACS) are downscaled using higher-resolution "
            "Census blocks, then further refined using city/county zoning data. "
            "People are placed only where residential parcels actually exist."
        ),
        source="Data: U.S. Census ACS 5-Year, Decennial Census, Orange County parcels",
    )

    # Slide 3: Travel-time pipeline
    slide_pipeline(
        prs,
        title="What is a School Community: Travel Time Analysis",
        images=[
            ss / "step1_road_network.png",
            ss / "step2_travel_time.png",
            ss / "step3_communities.png",
        ],
        labels=[
            "Road Network\n(speeds, signals, bike paths)",
            "Travel Time Calculation\n(Dijkstra shortest path)",
            "Parcels + Proximity\n(where people meet schools)",
        ],
        caption=(
            "Real road networks with speed limits, traffic signals, and bike "
            "infrastructure are used to calculate travel times from every parcel "
            "to each school. Communities are defined by actual travel-time "
            "proximity, not straight-line distance."
        ),
        source=(
            "Data: OpenStreetMap road network, NCDOT speed data, "
            "Overpass API intersection controls"
        ),
    )

    # Slide 4: Synthesis
    slide_synthesis(prs, ss)

    # Slide 5: Summary
    slide_summary(prs)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"\nPresentation saved to {output_path}")


# ── CLI ──────────────────────────────────────────────────────────────────────

def main():
    ap = argparse.ArgumentParser(
        description="Build CHCCS methods-overview slide deck")
    ap.add_argument("--skip-screenshots", action="store_true",
                    help="Reuse cached screenshots from a previous run")
    ap.add_argument("--output", type=Path, default=DEFAULT_OUTPUT,
                    help="Output .pptx path (default: assets/slides/methods_overview.pptx)")
    args = ap.parse_args()

    if not args.skip_screenshots:
        capture_all_screenshots(SCREENSHOTS_DIR)
    else:
        pngs = list(SCREENSHOTS_DIR.glob("*.png")) if SCREENSHOTS_DIR.exists() else []
        if not pngs:
            print("No cached screenshots found. Run without --skip-screenshots first.")
            return
        print(f"Using {len(pngs)} cached screenshots.")

    build_presentation(SCREENSHOTS_DIR, args.output)


if __name__ == "__main__":
    main()
